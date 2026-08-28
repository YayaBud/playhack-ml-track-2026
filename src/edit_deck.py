"""
Update PlayHack_ML_Round1_Submission.pptx to the corrected, metric-aware story.

The deck was built before the organisers' scoring rule was located, so it argues
for the F1-optimal threshold 0.290 -- which we now know scores 0.000 on both
Task B components. This rewrites the affected slides, inserts a new slide for the
metric itself, and swaps in the threshold-sweep figures.

Run:  python src/edit_deck.py
Out:  PlayHack_ML_Round1_Submission.pptx  (original backed up by the caller)
"""
from __future__ import annotations

import copy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "PlayHack_ML_Round1_Submission.pptx"
REPORTS = ROOT / "reports"

# palette lifted from the existing deck
DEEP = RGBColor(0x06, 0x5A, 0x82)
TEAL = RGBColor(0x1C, 0x72, 0x93)
MIDNIGHT = RGBColor(0x21, 0x29, 0x5C)
LIGHTFILL = RGBColor(0xEA, 0xF1, 0xF5)
MUTED = RGBColor(0x6B, 0x77, 0x85)
INK = RGBColor(0x1B, 0x1F, 0x27)
GREEN = RGBColor(0x6F, 0xCF, 0x97)
RED = RGBColor(0xD6, 0x45, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xC7, 0xD6, 0xE0)

HEAD, BODY, MONO = "Cambria", "Calibri", "Courier New"


# ------------------------------------------------------------- text helpers ---
def set_text(shape, text):
    """Replace a shape's text, keeping the first run's formatting."""
    tf = shape.text_frame
    para = tf.paragraphs[0]
    if not para.runs:
        para.add_run()
    para.runs[0].text = text
    for r in para.runs[1:]:
        r._r.getparent().remove(r._r)
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)


def find(slide, needle, exact=False):
    """First shape whose text matches."""
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if (t == needle) if exact else (needle.lower() in t.lower()):
                return sh
    return None


def retext(slide, needle, new, exact=False):
    sh = find(slide, needle, exact)
    if sh is None:
        print("    ! not found: %r" % needle[:50])
        return False
    set_text(sh, new)
    return True


def add_text(slide, x, y, w, h, text, size, *, bold=False, color=INK,
             font=BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    if spacing:
        r.font._rPr.set("spc", str(int(spacing * 100)))
    return box


def add_rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.04
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    if sh.has_text_frame:
        sh.text_frame.text = ""
    return sh


def add_image_fit(slide, path, x, y, w, h):
    """Insert an image centred inside the box, preserving aspect ratio."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    if w / h > ar:                 # box wider than image -> height-bound
        dh, dw = h, h * ar
    else:
        dw, dh = w, w / ar
    return slide.shapes.add_picture(str(path), Inches(x + (w - dw) / 2),
                                    Inches(y + (h - dh) / 2),
                                    Inches(dw), Inches(dh))


def move_slide(prs, old_index, new_index):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    lst.remove(ids[old_index])
    lst.insert(new_index, ids[old_index])


def main():
    prs = Presentation(str(DECK))
    S = prs.slides
    print("opened deck: %d slides" % len(S))

    # ------------------------------------------------------------ slide 1 ---
    print("\n[S1] title chips")
    s1 = S[0]
    retext(s1, "99.7% Bayes ceiling closed", "97% of metric ceiling closed")
    retext(s1, "AutoGluon SOTA blend", "Tuned to the official metric")
    retext(s1, "A leak-safe ML pipeline",
           "A leak-safe ML pipeline that turns 30 days of wearable data into injury "
           "risk, onset timing, and recovery-length predictions.")

    # ------------------------------------------------------------ slide 9 ---
    # Reframe the ceiling from ROC-AUC (never scored) to the official metric.
    print("[S9] ceiling -> official metric")
    s9 = S[8]
    retext(s9, "WHY 0.77 AND NOT 1.0", "WHY NOT A PERFECT SCORE")
    retext(s9, "We proved the ceiling",
           "We proved the ceiling — and closed 97% of it")
    retext(s9, "100 athletes, identical sensor data",
           "100 athletes, identical sensor data, true injury probability p = 0.70 → "
           "nature still flips a biased coin. Even a God-mode oracle cannot say which "
           "30 dodge it, so no model reaches a perfect score. We derived that limit "
           "for the metric we are actually scored on, then measured our distance to it.")
    retext(s9, "3 — Closed-form",
           "3 — Bayes-optimal F1 and skill scores under the official rule")
    retext(s9, "AUC* (Oracle)",
           "Metric ceiling 0.4407      Ours 0.4275      (ROC-AUC 0.7625 vs 0.7702, diagnostic)")
    retext(s9, "99.7%", "97.0%", exact=True)
    for sh in s9.shapes:
        if sh.has_table:
            rows = [("Task A — F1", "0.5210", "0.5246"),
                    ("Task B — onset skill", "0.6541", "0.6647"),
                    ("Task B — recovery skill", "0.1075", "0.1329")]
            for ri, vals in enumerate(rows, start=1):
                for ci, v in enumerate(vals):
                    cell = sh.table.cell(ri, ci)
                    if cell.text_frame.paragraphs[0].runs:
                        cell.text_frame.paragraphs[0].runs[0].text = v
                        for extra in cell.text_frame.paragraphs[0].runs[1:]:
                            extra._r.getparent().remove(extra._r)
            print("    table -> official metric")

    # ----------------------------------------------------------- slide 10 ---
    # Old S10 argued for threshold 0.290. Rewrite as the sweep evidence.
    print("[S10] threshold slide -> sweep evidence")
    s10 = S[9]
    retext(s10, "Why 0.290, not the default 0.5",
           "Why we flag 99.4%, not the F1 optimum")
    retext(s10, "Left peak = background hazard",
           "102 thresholds swept out-of-fold · one missed injury measurably dents "
           "recovery skill · 22 misses zero it")
    for sh in list(s10.shapes):
        if sh.has_chart:
            sh._element.getparent().remove(sh._element)
            print("    removed stale bimodal chart")
    # slide-optimised single-panel version; the two-panel document figure is
    # unreadable once scaled into this box
    add_image_fit(s10, REPORTS / "10_sweep_slide.png", 0.60, 1.62, 7.40, 4.65)

    cards = [("0.38", "F1-optimal", "Misses 502 of 1,050 injuries — both skills 0.000"),
             ("0.046", "Ours — shipped", "Flags 99.4%, misses zero, pays no penalty"),
             ("3rd", "of 102 thresholds", "Best or within 0.001 under 6 of 7 weightings")]
    olds = [("0.5", "Default", "Misses much of the overload cluster"),
            ("0.385", "Mean probability", "Floods false positives"),
            ("0.290", "F1-optimal — ours", "Sits in the gap")]
    for (ov, ol, od), (nv, nl, nd) in zip(olds, cards):
        retext(s10, ov, nv, exact=True)
        retext(s10, ol, nl)
        retext(s10, od, nd)

    # ------------------------------------------------- new slide: the metric ---
    print("[NEW] inserting 'the metric' slide after S9")
    blank = prs.slide_masters[0].slide_layouts[0]
    ns = S.add_slide(blank)
    # carry over the page-number chrome from S10 so the deck stays consistent
    for sh in S[9].shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "10":
            ns.shapes._spTree.append(copy.deepcopy(sh._element))

    add_text(ns, 0.60, 0.35, 8.0, 0.35, "THE SCORING RULE", 12, bold=True,
             color=TEAL, font=BODY, spacing=1.5)
    add_text(ns, 0.60, 0.65, 12.10, 0.80,
             "A missed injury costs 30. The baseline scores 3.2.", 30,
             bold=True, color=MIDNIGHT, font=HEAD)

    add_image_fit(ns, REPORTS / "08_penalty_cliff.png", 0.60, 1.65, 7.30, 3.45)

    add_text(ns, 0.60, 5.45, 7.30, 1.30,
             "Task B is scored only over athletes who are truly injured. A false "
             "positive is not in that population, so it costs Task B nothing — but a "
             "miss costs a flat 30 on both timing heads, against baselines of just "
             "7.61 and 3.24. One miss wipes out roughly nine good recovery "
             "predictions.", 12.5, color=INK, font=BODY)

    box = add_rect(ns, 8.25, 1.65, 4.45, 2.05, LIGHTFILL)
    add_text(ns, 8.55, 1.90, 3.90, 0.35, "BREAK-EVEN RECALL", 11, bold=True,
             color=TEAL, font=BODY, spacing=1.2)
    # keep the number boxes narrow so they cannot run into their labels
    add_text(ns, 8.55, 2.30, 1.55, 0.45, "0.82", 22, bold=True, color=MIDNIGHT, font=HEAD)
    add_text(ns, 10.25, 2.42, 2.20, 0.30, "onset skill", 11.5, color=MUTED, font=BODY)
    add_text(ns, 8.55, 2.85, 1.55, 0.45, "0.99", 22, bold=True, color=RED, font=HEAD)
    add_text(ns, 10.25, 2.97, 2.20, 0.30, "recovery skill", 11.5, color=MUTED, font=BODY)

    # Side-by-side comparison rather than one big "0.000": a lone zero on a dark
    # card reads at a glance as OUR score, which is the opposite of the point.
    box2 = add_rect(ns, 8.25, 3.95, 4.45, 2.80, MIDNIGHT)
    add_text(ns, 8.55, 4.20, 3.90, 0.35, "TASK B SCORE, BY THRESHOLD", 11, bold=True,
             color=GREEN, font=BODY, spacing=1.2)

    add_text(ns, 8.55, 4.68, 2.35, 0.30, "Tuned for F1", 12.5, color=PALE, font=BODY)
    add_text(ns, 11.00, 4.60, 1.45, 0.45, "0.000", 20, bold=True,
             color=RED, font=HEAD, align=PP_ALIGN.RIGHT)

    add_text(ns, 8.55, 5.28, 2.35, 0.30, "Our threshold", 12.5, bold=True,
             color=WHITE, font=BODY)
    add_text(ns, 11.00, 5.20, 1.45, 0.45, "0.762", 20, bold=True,
             color=GREEN, font=HEAD, align=PP_ALIGN.RIGHT)

    add_text(ns, 8.55, 5.85, 3.90, 0.75,
             "Same models, same predictions — only the cut-off differs. It roughly "
             "doubles the combined score.", 11.5, color=PALE, font=BODY)

    move_slide(prs, len(S._sldIdLst) - 1, 9)
    print("    placed at position 10")

    # ----------------------------------------------------------- results ----
    print("[S12] results table -> official metric")
    s_res = S[11]
    for sh in s_res.shapes:
        if sh.has_table:
            t = sh.table
            hdr = ["Task", "Scored quantity", "Baseline", "Ours", "Ceiling", "Closed"]
            rows = [("Task A", "F1", "—", "0.5210", "0.5246", "99.3%"),
                    ("Task B", "skill — onset", "0.000", "0.6541", "0.6647", "98.4%"),
                    ("Task B", "skill — recovery", "0.000", "0.1075", "0.1329", "80.9%"),
                    ("Combined", "mean of three", "—", "0.4275", "0.4407", "97.0%")]
            for ci, v in enumerate(hdr):
                c = t.cell(0, ci)
                if c.text_frame.paragraphs[0].runs:
                    c.text_frame.paragraphs[0].runs[0].text = v
                    for e in c.text_frame.paragraphs[0].runs[1:]:
                        e._r.getparent().remove(e._r)
            for ri, vals in enumerate(rows, start=1):
                for ci, v in enumerate(vals):
                    c = t.cell(ri, ci)
                    if c.text_frame.paragraphs[0].runs:
                        c.text_frame.paragraphs[0].runs[0].text = v
                        for e in c.text_frame.paragraphs[0].runs[1:]:
                            e._r.getparent().remove(e._r)
            print("    table updated")
    retext(s_res, "Decision threshold: 0.290",
           "Decision threshold: 0.046 (rank-quantile)  ·  Test flagged: 99.4%  ·  "
           "OOF recall on injured: 1.000  ·  Features per head: 35 / 35 / 20")
    retext(s_res, "Custom 3-model ensemble alone",
           "ROC-AUC 0.7625 (ceiling 0.7702) is a diagnostic, not a scored quantity. "
           "F1 of 0.52 is not a weak classifier: at full recall F1 is pinned by "
           "prevalence to 2p/(1+p) = 0.5185, and the Bayes optimum under this metric "
           "is 0.5226.")

    # -------------------------------------------------- negative results ----
    print("[S13] negative results: remove stale figures")
    s_neg = S[12]
    retext(s_neg, "Lost on both heads",
           "Lost on both heads against the tuned ensemble. A single parametric hazard "
           "curve cannot represent two distinct biological mechanisms.")
    retext(s_neg, "206 features scored below",
           "206 features scored 0.7497 AUC — below a single raw hr_acwr column "
           "(0.7583). The ACWR family's r ≈ 0.99 redundancy makes trees overfit "
           "near-identical splits. 35 features scored best, at 0.7572.")

    # ---------------------------------------------------------- takeaways ---
    print("[S14] takeaways")
    s_take = S[13]
    retext(s_take, "A closed-form Bayes-optimal AUC",
           "We derived the Bayes optimum for the organisers' own rule — and from the "
           "data generating process independently — to show 0.52 F1 is 99.3% of what "
           "is achievable, not a weak model.")

    # ------------------------------------------------------- housekeeping ---
    # Inserting a slide leaves every page number after it off by one.
    print("[all] renumbering pages")
    from pptx.util import Emu as _Emu
    fixed = 0
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text.strip()
            try:
                x, y = _Emu(sh.left).inches, _Emu(sh.top).inches
            except Exception:
                continue
            if x > 11.5 and y > 6.5 and len(t) <= 3 and t.isdigit():
                want = "%02d" % i
                if t != want:
                    set_text(sh, want)
                    fixed += 1
    print("    corrected %d page numbers" % fixed)

    # team name removed at the user's request
    for sh in list(prs.slides[0].shapes):
        if sh.has_text_frame and "yayabud" in sh.text_frame.text.lower():
            sh._element.getparent().remove(sh._element)
            print("[S1] removed team name")

    # our own threshold should not be flagged in alarm-red on the sweep slide
    for sh in prs.slides[10].shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "0.046":
            for p_ in sh.text_frame.paragraphs:
                for r in p_.runs:
                    r.font.color.rgb = TEAL
            print("[S11] recoloured 0.046 from red to teal")

    prs.save(str(DECK))
    print("\nsaved %s  (%d slides)" % (DECK.name, len(prs.slides)))


if __name__ == "__main__":
    main()
